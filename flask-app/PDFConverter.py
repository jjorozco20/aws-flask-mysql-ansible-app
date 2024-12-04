import os
from pdf2image import convert_from_path
from pptx import Presentation
from pptx.util import Inches

# Directory containing the PDF files
input_folder = "C:\\Users\\Juan.Orozco\\Downloads\\slides"  # Replace with your folder path
output_folder = "C:\\Poppler"  # Replace with your desired output path

# Ensure the output folder exists
os.makedirs(output_folder, exist_ok=True)

# Path to the Poppler binaries (ensure this is correct for your system)
poppler_path = r'C:\ProgramData\chocolatey\bin'

# Step 1: Loop through each PDF file in the input folder
for filename in os.listdir(input_folder):
    if filename.lower().endswith('.pdf'):
        pdf_path = os.path.join(input_folder, filename)
        print(f"Processing: {filename}")

        # Step 2: Convert PDF pages to images with the Poppler path specified
        pages = convert_from_path(pdf_path, dpi=300, poppler_path=poppler_path)  # High resolution for clarity

        # Step 3: Create a new PowerPoint presentation
        presentation = Presentation()

        # Step 4: Add each PDF page as an image to the slides
        for i, page in enumerate(pages):
            # Save the page as an image file
            temp_image = f"temp_slide_{i+1}.png"
            page.save(temp_image, "PNG")

            # Create a blank slide
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])  # Layout 6 is a blank slide

            # Add the image to the slide
            left = Inches(0)
            top = Inches(0)
            slide.shapes.add_picture(temp_image, left, top, width=Inches(10), height=Inches(7.5))

            # Remove the temporary image after use
            os.remove(temp_image)

        # Step 5: Replace "Caleb Espinoza" with "Juan Orozco" in text boxes on each slide
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    if "Caleb Espinoza" in shape.text:
                        # Replace text in the shape
                        shape.text = shape.text.replace("Caleb Espinoza", "Juan Orozco")

        # Step 6: Save the presentation
        output_pptx_path = os.path.join(output_folder, f"{os.path.splitext(filename)[0]}.pptx")
        presentation.save(output_pptx_path)
        print(f"PowerPoint file created: {output_pptx_path}")

print("All PDF files processed and text replaced successfully!")
